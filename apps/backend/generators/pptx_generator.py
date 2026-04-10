"""Generador de presentaciones PowerPoint a partir del análisis de IA."""
import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from core.models import AnalysisResult
from core.session import Session

logger = logging.getLogger(__name__)

# Paleta corporativa
CLR_DARK = RGBColor(0x1F, 0x49, 0x7D)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_ACCENT = RGBColor(0x2E, 0x75, 0xB6)
CLR_BODY = RGBColor(0x26, 0x26, 0x26)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def generate_pptx(result: AnalysisResult, session: Session, template_path: str | None = None) -> Path:
    """
    Genera un archivo PPTX a partir del análisis de IA.
    El texto debe contener bloques separados por '---SLIDE---'.
    Si se proporciona template_path, usa esa plantilla .pptx como base.
    Retorna la ruta del archivo en el directorio de sesión.
    """
    # Extraer la sección de slides si existe el separador ===SLIDES===
    raw = result.analysis_text
    if "===SLIDES===" in raw:
        raw = raw.split("===SLIDES===")[1].strip()

    slide_blocks = _parse_slide_blocks(raw)
    if not slide_blocks:
        # Si no hay bloques estructurados, generar slides desde el Markdown
        slide_blocks = _fallback_parse_markdown(raw)

    # Usar plantilla personalizada si se proporcionó y existe
    if template_path:
        try:
            prs = Presentation(template_path)
            # Eliminar todas las slides existentes de la plantilla preservando el master/layouts
            from pptx.oxml.ns import qn
            sldIdLst = prs.slides._sldIdLst
            for sldId in list(sldIdLst):
                r_id = sldId.get(qn('r:id'))
                prs.part.drop_rel(r_id)
                sldIdLst.remove(sldId)
            logger.info("Usando plantilla personalizada: %s", template_path)
        except Exception as exc:
            logger.warning("No se pudo cargar la plantilla '%s': %s — usando plantilla vacía", template_path, exc)
            prs = Presentation()
    else:
        prs = Presentation()

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, block in enumerate(slide_blocks):
        slide_type = block.get("tipo", "contenido").lower()
        if i == 0 or slide_type == "portada":
            _add_cover_slide(prs, block)
        elif slide_type == "cierre":
            _add_closing_slide(prs, block)
        else:
            _add_content_slide(prs, block)

    out_path = session.temp_dir / "presentacion_analisis.pptx"
    prs.save(str(out_path))
    logger.info("Presentación PPTX generada: %s", out_path)
    return out_path


# ────────────────────────────────────────────────────────────────────────────
# Parsers
# ────────────────────────────────────────────────────────────────────────────

def _parse_slide_blocks(text: str) -> list[dict]:
    """Parsea bloques delimitados por ---SLIDE--- o === SLIDE N: titulo ===."""
    # Normalizar separadores alternativos que algunos LLMs generan
    # Convierte "=== SLIDE 1: Titulo ===" en "---SLIDE---\nTITULO: Titulo"
    normalized = re.sub(
        r"===+\s*SLIDE\s*\d*:?\s*([^=]*?)\s*===+",
        lambda m: f"---SLIDE---\nTITULO: {m.group(1).strip()}",
        text,
        flags=re.IGNORECASE,
    )

    blocks = []
    raw_blocks = re.split(r"---SLIDE---", normalized, flags=re.IGNORECASE)
    for raw in raw_blocks:
        raw = raw.strip()
        if not raw:
            continue
        block: dict = {"tipo": "contenido", "titulo": "", "bullets": []}
        for line in raw.splitlines():
            line = line.strip()
            if line.upper().startswith("TIPO:"):
                block["tipo"] = line.split(":", 1)[1].strip()
            elif line.upper().startswith("TITULO:"):
                block["titulo"] = line.split(":", 1)[1].strip()
            elif line.startswith("- ") or line.startswith("* "):
                # Limpiar markdown bold (**texto**) de los bullets
                bullet = re.sub(r"\*\*(.+?)\*\*", r"\1", line[2:].strip())
                block["bullets"].append(bullet)
        if block["titulo"] or block["bullets"]:
            blocks.append(block)
    return blocks


def _fallback_parse_markdown(text: str) -> list[dict]:
    """Convierte Markdown en slides agrupando bajo encabezados ##."""
    slides: list[dict] = []
    current: dict | None = None
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("## ") or stripped.startswith("# "):
            if current and (current["titulo"] or current["bullets"]):
                slides.append(current)
            current = {"tipo": "contenido", "titulo": stripped.lstrip("#").strip(), "bullets": []}
        elif stripped.startswith(("- ", "* ")) and current is not None:
            current["bullets"].append(stripped[2:].strip())
        elif stripped and current is not None and not stripped.startswith("#"):
            # Línea de texto normal: añadir como bullet si es corta
            if len(stripped) <= 120:
                current["bullets"].append(stripped)
    if current and (current["titulo"] or current["bullets"]):
        slides.append(current)
    if not slides:
        slides = [{"tipo": "contenido", "titulo": "Análisis", "bullets": [text[:200]]}]
    return slides


# ────────────────────────────────────────────────────────────────────────────
# Constructores de slides
# ────────────────────────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(blank_layout)


def _add_cover_slide(prs: Presentation, block: dict) -> None:
    slide = _blank_slide(prs)
    # Fondo azul
    _fill_background(slide, CLR_DARK)
    # Título centrado
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = block.get("titulo", "Análisis Ejecutivo")
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = CLR_WHITE

    # Subtítulo / primer bullet
    bullets = block.get("bullets", [])
    if bullets:
        txSub = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.33), Inches(1))
        tf2 = txSub.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = bullets[0]
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)


def _add_content_slide(prs: Presentation, block: dict) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, CLR_WHITE)

    # Banda superior azul
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(1.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_DARK
    bar.line.fill.background()

    # Título en la banda
    txTitle = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.9))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = block.get("titulo", "")
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = CLR_WHITE

    # Bullets
    bullets = block.get("bullets", [])[:6]
    if bullets:
        txBody = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8))
        tf2 = txBody.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(18)
            run.font.color.rgb = CLR_BODY
            p.space_after = Pt(8)


def _add_closing_slide(prs: Presentation, block: dict) -> None:
    slide = _blank_slide(prs)
    _fill_background(slide, CLR_ACCENT)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.33), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = block.get("titulo", "Conclusiones")
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = CLR_WHITE

    bullets = block.get("bullets", [])
    if bullets:
        txSub = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(10.33), Inches(2))
        tf2 = txSub.text_frame
        tf2.word_wrap = True
        for i, b in enumerate(bullets[:3]):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r = p2.add_run()
            r.text = b
            r.font.size = Pt(16)
            r.font.color.rgb = CLR_WHITE


def _fill_background(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color
